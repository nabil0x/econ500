#!/usr/bin/env python3
"""Extract text from PDF/DOCX/PPTX/IMG files into markdown for MkDocs."""

import os, sys, subprocess, tempfile, shutil, re
from pathlib import Path

# Optional import wrappers
HAS_PDFMINER = False
HAS_PDFPLUMBER = False
HAS_PYDOCX = False
HAS_PPPPTX = False
HAS_PYTESS = False
HAS_PIL = False

try:
    from pdfminer.high_level import extract_text as pdfminer_extract
    HAS_PDFMINER = True
except ImportError:
    pass

try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    pass

try:
    import docx
    HAS_PYDOCX = True
except ImportError:
    pass

try:
    from pptx import Presentation
    HAS_PPPPTX = True
except ImportError:
    pass

try:
    import pytesseract
    HAS_PYTESS = True
except ImportError:
    pass

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    pass

REPO_ROOT = Path(__file__).resolve().parent.parent
EXTRACTED = REPO_ROOT / "_support" / "extracted"

SKIP_EXTS = {".md", ".html", ".htm", ".css", ".js", ".yml", ".yaml",
             ".json", ".svg", ".py", ".gitignore", ".gitattributes",
             ".map", ".zip", ".png", ".jpg", ".jpeg", ".ico"}
SKIP_DIRS = {".git", "node_modules", "__pycache__", ".opencode", ".omo"}
SKIP_FILES = {"mkdocs.yml", "requirements.txt", "README.md"}

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"}

def slugify(name: str) -> str:
    """Create a safe filename from a title."""
    name = re.sub(r'[^\w\s-]', '', name)
    name = re.sub(r'[-\s]+', '-', name)
    return name.strip('-').lower() or 'untitled'

def extract_pdf_pdftotext(path: Path) -> str | None:
    """Extract text via pdftotext (poppler). Fast path for text PDFs."""
    try:
        result = subprocess.run(
            ["pdftotext", "-layout", "-nopgbrk", str(path), "-"],
            capture_output=True, text=True, timeout=120
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout
    except (subprocess.TimeoutExpired, FileNotFoundError):
        pass
    return None

def extract_pdf_pdfminer(path: Path) -> str | None:
    """Extract text via pdfminer.six."""
    if not HAS_PDFMINER:
        return None
    try:
        text = pdfminer_extract(str(path))
        if text and text.strip():
            return text
    except Exception:
        pass
    return None

def extract_pdf_pdfplumber(path: Path) -> str | None:
    """Extract text via pdfplumber."""
    if not HAS_PDFPLUMBER:
        return None
    try:
        with pdfplumber.open(str(path)) as pdf:
            pages = []
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    pages.append(t)
            if pages:
                return "\n\n".join(pages)
    except Exception:
        pass
    return None

def extract_pdf_ocr(path: Path) -> str | None:
    """OCR a PDF page by page using tesseract. Requires pdftoppm."""
    if not HAS_PYTESS or not HAS_PIL:
        return None
    if not shutil.which("pdftoppm"):
        return None

    with tempfile.TemporaryDirectory() as tmpd:
        try:
            subprocess.run(
                ["pdftoppm", "-png", "-r", "300", str(path), f"{tmpd}/page"],
                capture_output=True, timeout=300
            )
            pages = sorted(Path(tmpd).glob("page-*.png"))
            if not pages:
                return None
            parts = []
            for pg in pages:
                try:
                    img = Image.open(pg)
                    text = pytesseract.image_to_string(img, lang='eng')
                    parts.append(text)
                except Exception:
                    parts.append("")
            return "\n\n".join(parts)
        except (subprocess.TimeoutExpired, Exception):
            return None

def extract_pdf(path: Path) -> str | None:
    """Extract text from a PDF. Returns markdown-safe text."""
    # Fast path: pdftotext
    text = extract_pdf_pdftotext(path)
    if text and len(text.strip()) > 100:
        return text

    # Second: pdfminer
    text = extract_pdf_pdfminer(path)
    if text and len(text.strip()) > 50:
        return text

    # Third: pdfplumber
    text = extract_pdf_pdfplumber(path)
    if text and len(text.strip()) > 50:
        return text

    # Fallback: OCR (slow, only if we have few chars)
    print(f"  [OCR] {path.name}")
    return extract_pdf_ocr(path)

def extract_docx(path: Path) -> str | None:
    """Convert DOCX to markdown text."""
    if not HAS_PYDOCX:
        return None
    try:
        doc = docx.Document(str(path))
        lines = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                style = para.style.name.lower() if para.style else ""
                if "heading 1" in style:
                    lines.append(f"# {text}")
                elif "heading 2" in style:
                    lines.append(f"## {text}")
                elif "heading 3" in style:
                    lines.append(f"### {text}")
                else:
                    lines.append(text)
        return "\n\n".join(lines)
    except Exception:
        return None

def extract_pptx(path: Path) -> str | None:
    """Convert PPTX to markdown text."""
    if not HAS_PPPPTX:
        return None
    try:
        prs = Presentation(str(path))
        pages = []
        for i, slide in enumerate(prs.slides, 1):
            lines = [f"## Slide {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lines.append(t)
            pages.append("\n\n".join(lines))
        return "\n\n---\n\n".join(pages)
    except Exception:
        return None

def extract_image(path: Path) -> str | None:
    """OCR an image file."""
    if not HAS_PYTESS or not HAS_PIL:
        return None
    try:
        img = Image.open(path)
        text = pytesseract.image_to_string(img, lang='eng')
        if text and text.strip():
            return text
    except Exception:
        pass
    return None

def get_output_path(src: Path) -> Path:
    """Determine output md path mirroring repo structure."""
    rel = src.relative_to(REPO_ROOT)
    # Replace extension with .md, remove special chars
    stem = src.stem
    # Clean the filename for mkdocs
    out_name = slugify(stem) + ".md"
    parent = rel.parent
    return EXTRACTED / parent / out_name

def needs_extraction(path: Path) -> bool:
    """Check if file needs text extraction."""
    ext = path.suffix.lower()
    if ext not in {".pdf", ".docx", ".pptx", ".ppt"} and ext not in IMAGE_EXTS:
        return False
    # Skip files over 50MB (likely scanned textbooks that OCR can't handle)
    try:
        if path.stat().st_size > 50 * 1024 * 1024:
            print(f"  SKIP (>{50}MB)")
            return False
    except OSError:
        pass
    return True

def process_file(path: Path) -> bool:
    """Extract text from file and write markdown. Returns True on success."""
    ext = path.suffix.lower()
    out_path = get_output_path(path)
    if out_path.exists():
        return True  # already extracted

    out_path.parent.mkdir(parents=True, exist_ok=True)

    text = None
    if ext == ".pdf":
        text = extract_pdf(path)
    elif ext == ".docx":
        text = extract_docx(path)
    elif ext in {".pptx", ".ppt"}:
        text = extract_pptx(path)
    elif ext in IMAGE_EXTS:
        text = extract_image(path)

    if text and text.strip():
        lines = [
            f"<!-- Extracted from: {path.relative_to(REPO_ROOT)} -->",
            "",
            text.strip(),
            ""
        ]
        out_path.write_text("\n".join(lines), encoding="utf-8")
        print(f"  OK  {out_path.relative_to(EXTRACTED)}")
        return True
    else:
        # Write empty placeholder so we don't retry
        lines = [
            f"<!-- Extracted from: {path.relative_to(REPO_ROOT)} -->",
            "",
            "*No text could be extracted from this file.*",
            ""
        ]
        out_path.write_text("\n".join(lines), encoding="utf-8")
        print(f"  EMPTY {out_path.relative_to(EXTRACTED)}")
        return False

def scan_and_extract(target_dir: str | None = None):
    """Scan and extract text from all binary files."""
    base = REPO_ROOT / target_dir if target_dir else REPO_ROOT
    if not base.exists():
        print(f"Directory not found: {base}")
        return

    count = 0
    success = 0
    for root, dirs, files in os.walk(base):
        # Skip unwanted dirs
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS and not d.startswith(".")]

        for fname in sorted(files):
            fpath = Path(root) / fname
            if not needs_extraction(fpath):
                continue
            if fname in SKIP_FILES:
                continue

            count += 1
            sys.stdout.write(f"[{count}] {fpath.relative_to(REPO_ROOT)} ... ")
            sys.stdout.flush()
            if process_file(fpath):
                success += 1

    print(f"\nDone: {success}/{count} files extracted successfully.")

if __name__ == "__main__":
    target = sys.argv[1] if len(sys.argv) > 1 else None
    scan_and_extract(target)
